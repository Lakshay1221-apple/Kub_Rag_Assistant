"""Office document loader for DOCX and PPTX files."""

from __future__ import annotations

from pathlib import Path

import logfire
from docx import Document
from loguru import logger
from pptx import Presentation

__all__ = ["parse_office"]


def _validate_file_path(file_path: str) -> Path:
	"""Validate that the provided path exists and points to a file."""

	path = Path(file_path)
	if not path.exists():
		logger.error("Office file does not exist: {}", file_path)
		raise FileNotFoundError(f"Office file does not exist: {file_path}")
	if not path.is_file():
		logger.error("Office path is not a file: {}", file_path)
		raise IsADirectoryError(f"Office path is not a file: {file_path}")
	return path


def _normalize_text(text: str) -> str:
	"""Normalize whitespace while preserving readable line breaks."""

	return "\n".join(line.strip() for line in text.splitlines() if line.strip())


def _extract_docx_text(path: Path) -> str:
	"""Extract all non-empty paragraph text from a DOCX file."""

	logger.info("Extracting DOCX content: {}", path)
	document = Document(str(path))
	paragraphs = [paragraph.text.strip() for paragraph in document.paragraphs if paragraph.text.strip()]
	return "\n".join(paragraphs)


def _extract_pptx_text(path: Path) -> str:
	"""Extract text from all shapes across all slides in a PPTX file."""

	logger.info("Extracting PPTX content: {}", path)
	presentation = Presentation(str(path))
	slide_text_parts: list[str] = []

	for slide_index, slide in enumerate(presentation.slides, start=1):
		slide_parts: list[str] = []
		for shape in slide.shapes:
			if not hasattr(shape, "text"):
				continue
			text = getattr(shape, "text", "")
			if text and text.strip():
				slide_parts.append(text.strip())

		if slide_parts:
			slide_text_parts.append("\n".join(slide_parts))
			logger.debug("Extracted text from slide {} in {}", slide_index, path)
		else:
			logger.debug("Skipping empty slide {} in {}", slide_index, path)

	return "\n\n".join(slide_text_parts)


def parse_office(file_path: str) -> str:
	"""Parse DOCX or PPTX content and return cleaned text."""

	logger.info("Starting office parse: {}", file_path)
	try:
		path = _validate_file_path(file_path)
		suffix = path.suffix.lower()

		if suffix == ".docx":
			raw_text = _extract_docx_text(path)
		elif suffix == ".pptx":
			raw_text = _extract_pptx_text(path)
		else:
			logger.error("Unsupported office file type: {}", file_path)
			raise ValueError(f"Unsupported office file type: {file_path}")

		clean_text = _normalize_text(raw_text)
		logger.info("Completed office parse: {} characters extracted", len(clean_text))
		return clean_text
	except Exception as exc:
		logfire.error(f"Office parse failed: {exc}")
		logger.exception("Office parse failed for {}", file_path)
		raise
